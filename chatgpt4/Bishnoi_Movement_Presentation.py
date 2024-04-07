
# https://chat.openai.com/c/d2c762af-856a-4b83-8b5c-cf8b218b3610

from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
prs = Presentation()

# Slide Titles and Content
slides_content = [
    ("The Bishnoi Movement: Pioneers of Environmental Conservation", "A Historical Insight into Sustainable Living\nYour Name\nDate\nCollege/University Name"),
    ("Introduction to the Bishnoi Movement", "Overview of the movement\nOrigin: 15th century, Rajasthan, India\nFounder: Guru Jambheshwar (Jambhoji)"),
    ("The 29 Principles", "Overview of the 29 principles\nFocus on environmental conservation and animal protection"),
    ("Cultural and Environmental Significance", "Dedication to nature as part of religious belief\nRole in preserving local flora and fauna"),
    ("The Khejarli Massacre: A Turning Point", "Story of Amrita Devi Bishnoi and the 363 sacrifices\nImpact on Indian environmental consciousness"),
    ("Legacy and Influence", "Inspiration for subsequent movements like Chipko\nInfluence on modern environmentalism"),
    ("The Bishnoi Community Today", "Current practices and conservation efforts\nLegal battles and activism"),
    ("Challenges Faced", "Modern challenges\nConflicts with development projects and poaching"),
    ("Global Relevance", "Universal appeal of Bishnoi principles\nLessons for sustainable living"),
    ("Conclusion and Reflection", "Importance of the Bishnoi movement\nApplying principles today for a sustainable future\nCall to action for conservation"),
]

# Add slides and content
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    body_placeholder.text = content

# Save the presentation
pptx_file = "/mnt/data/Bishnoi_Movement_Presentation.pptx"
prs.save(pptx_file)

print(pptx_file)
