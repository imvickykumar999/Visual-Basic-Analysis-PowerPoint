
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
# pip install python-pptx

from pptx import Presentation

path_to_presentation = 'static/PlayFair Ciphertext Encryption.pptx'
prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                
for i in text_runs:
    print(i, end='\n'+'-'*20+'\n')
