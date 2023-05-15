
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

def extract():
    # pip install aspose.slides
    import aspose.slides as slides

    # Load a presentation
    with slides.Presentation("PPTs/presentation.ppt") as presentation:
        
        # Check if presentation contains VBA Project
        if presentation.vba_project is not None:
            
            # Print each module
            for module in presentation.vba_project.modules:
                print(module.name)
                print(module.source_code)
        else:
            print('Presentation does NOT contains VBA Project')


def read(path_to_presentation):
    # pip install python-pptx

    from pptx import Presentation
    prs = Presentation(path_to_presentation)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    # print(list(prs.slides))

    for slide in prs.slides:
        for shape in slide.shapes:
            print(shape.textframe.paragraphs)
        print()
            # if not shape.has_textframe:
            #     continue
            # for paragraph in shape.textframe.paragraphs:
            #     for run in paragraph.runs:
            #         text_runs.append(run.text)

path_to_presentation = 'PPTs/PlayFair Ciphertext Encryption.pptx'
read(path_to_presentation)
